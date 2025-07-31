# app.py

import os
import sys
import json
import pandas as pd
import io
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
def resource_path(rel_path: str) -> str:
    """Get absolute path to resource, works for dev and for PyInstaller."""
    if getattr(sys, "frozen", False):
        base = sys._MEIPASS
    else:
        base = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(base, rel_path)

# ─────────────────────────────────────────────────────────────────────────────
# Constants & Persistence
# ─────────────────────────────────────────────────────────────────────────────
BATCHES_PATH = "dashboards/batches.json"

def load_batches(path: str = BATCHES_PATH) -> list:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return []

def save_batches(batches: list, path: str = BATCHES_PATH) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(batches, f, indent=2, default=str)

# ─────────────────────────────────────────────────────────────────────────────
# Data Loading
# ─────────────────────────────────────────────────────────────────────────────
def load_dataframe(src) -> pd.DataFrame:
    # Handles file upload object or file path
    if hasattr(src, "read") and hasattr(src, "name"):
        data = src.getvalue()
        ext  = os.path.splitext(src.name)[1].lower()
        if ext == ".csv":
            return pd.read_csv(io.BytesIO(data), encoding="utf-8", engine="python", on_bad_lines="skip")
        elif ext in (".xls", ".xlsx"):
            return pd.read_excel(io.BytesIO(data))
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    ext = os.path.splitext(src)[1].lower()
    if ext == ".csv":
        return pd.read_csv(src)
    elif ext in (".xls", ".xlsx"):
        return pd.read_excel(src)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

# ─────────────────────────────────────────────────────────────────────────────
# Proposed Metrics Extraction
# ─────────────────────────────────────────────────────────────────────────────
def extract_proposed_metrics_anywhere(df):
    """
    Find 'Proposed Metrics' anywhere in the sheet, then extract the next 3 rows
    for 'Impressions', 'Engagements', 'Influencers' in the same column.
    Returns a dict: {'Impressions': ..., 'Engagements': ..., 'Influencers': ...}
    """
    found = False
    for col in df.columns:
        col_series = df[col].astype(str).str.strip().str.lower()
        idxs = col_series[col_series == "proposed metrics"].index.tolist()
        if idxs:
            found = True
            idx = idxs[0]
            col_num = df.columns.get_loc(col)
            names = []
            values = []
            for offset in range(1, 4):
                name = str(df.iloc[idx+offset, col_num]).strip()
                value = df.iloc[idx+offset, col_num+1]
                names.append(name)
                values.append(value)
            break
    if not found:
        raise ValueError("'Proposed Metrics' not found in any column.")
    return dict(zip(names, values))

# ─────────────────────────────────────────────────────────────────────────────
def format_compact_number(n):
    try:
        n = float(n)
        if n >= 1_000_000:
            return f"{n / 1_000_000:.1f}MM"
        elif n >= 1_000:
            return f"{n / 1_000:.1f}K"
        else:
            return str(int(n))
    except:
        return str(n)

# PowerPoint Deck Generation
# ─────────────────────────────────────────────────────────────────────────────
def populate_pptx_from_excel(excel_df, pptx_template_path, output_path, images=None, text_inputs=None):
    prs = Presentation(pptx_template_path)
    handle_slide_6 = text_inputs.get("slide_6", "@default")
    handle_slide_7_left = text_inputs.get("slide_7_left", "@default")
    handle_slide_7_right = text_inputs.get("slide_7_right", "@default")
    slide_7_likes = text_inputs.get("slide_7_like", "@default")
    slide_7_comments = text_inputs.get("slide_7_comment", "@default")
    slide_7_views = text_inputs.get("slide_7_view", "@default")
    slide_7_reach =  text_inputs.get("slide_7_reaches", "@default")
    slide_7_engagements = text_inputs.get("slide_7_eng", "@default")
    slide_7_impressions = text_inputs.get("slide_7_impr", "@default")
    text_slide_9 = text_inputs.get("slide_9", "@default")
    text_slide_13 = text_inputs.get("slide_13", "@default")
    text_slide_15 = text_inputs.get("slide_15", "@default")
    text_slide_16 = text_inputs.get("slide_16", "@default")
    date_slide_1 = text_inputs.get("slide_1_d", "@default")
    hashtag_slide_1 = text_inputs.get("slide_1_htg", "@default")
    date_slide_2 = text_inputs.get("slide_2_d", "@default")
    hashtag_slide_2 = text_inputs.get("slide_2_htg", "@default")
    date_slide_3 = text_inputs.get("slide_3_d", "@default")
    hashtag_slide_3 = text_inputs.get("slide_3_htg", "@default")
    slide_4_bullet1 = text_inputs.get("slide_4_b1", "@default")
    slide_4_bullet2 = text_inputs.get("slide_4_b2", "@default")

    # ---------- Extract Proposed Metrics Block (TextBox 2) ----------
    try:
        metrics = extract_proposed_metrics_anywhere(excel_df)
    except Exception as e:
        metrics = {"Impressions": "", "Engagements": "", "Influencers": ""}
        print(f"Warning: Could not extract Proposed Metrics from Excel: {e}")

    print("COLUMNS:", list(excel_df.columns))

    if images and "slide_6" in images and images["slide_6"] is not None:
        slide = prs.slides[5]
        img_bytes = images["slide_6"].read()
        temp_img_path = "temp_slide_6_img.jpg"
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)

    # Read the image size in pixels
        img = Image.open(temp_img_path)
        img_width_px, img_height_px = img.size

    # PowerPoint uses EMUs; there are 914400 EMUs per inch
    # Typical image DPI is 96, so we convert pixels to inches
        dpi = 96
        img_width = int(img_width_px / dpi * 914400)
        img_height = int(img_height_px / dpi * 914400)

        for shape in slide.shapes:
            if shape.name == "Picture 2":
                left, top = shape.left, shape.top
                box_width, box_height = shape.width, shape.height

            # --- Resize if necessary ---
                scale = min(
                    box_width / img_width if img_width > box_width else 1.0,
                    box_height / img_height if img_height > box_height else 1.0
                )
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)

            # Optional: Center inside the box if smaller than box
                final_left = left + int((box_width - final_width) / 2)
                final_top = top + int((box_height - final_height) / 2)

            # Remove the old placeholder/image
                slide.shapes._spTree.remove(shape._element)

            # Add image, resized if needed, centered in box
                slide.shapes.add_picture(temp_img_path, final_left, final_top, width=final_width, height=final_height)
                break


    # Slide 7
    slide = prs.slides[6]
    if images and "slide_7_left" in images and images["slide_7_left"] is not None:
        img_bytes = images["slide_7_left"].read()
        temp_img_path = "temp_slide_7_left.jpg"
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)

    # Get natural image size in pixels
        img = Image.open(temp_img_path)
        img_width_px, img_height_px = img.size

    # PowerPoint units (EMU)
        dpi = 96
        img_width = int(img_width_px / dpi * 914400)
        img_height = int(img_height_px / dpi * 914400)

        for shape in slide.shapes:
            if shape.name == "Picture 3":
                box_left, box_top = shape.left, shape.top
                box_width, box_height = shape.width, shape.height

            # Scale if needed
                scale = min(
                box_width / img_width if img_width > box_width else 1.0,
                box_height / img_height if img_height > box_height else 1.0
                )
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)

            # Center
                final_left = box_left + int((box_width - final_width) / 2)
                final_top = box_top + int((box_height - final_height) / 2)

            # Remove old placeholder
                slide.shapes._spTree.remove(shape._element)

            # Add the image, centered & fitted
                slide.shapes.add_picture(temp_img_path, final_left, final_top, width=final_width, height=final_height)
                break
    if images and "slide_7_right" in images and images["slide_7_right"] is not None:
        img_bytes = images["slide_7_right"].read()
        temp_img_path = "temp_slide_7_right.jpg"
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)

        img = Image.open(temp_img_path)
        img_width_px, img_height_px = img.size

        dpi = 96
        img_width = int(img_width_px / dpi * 914400)
        img_height = int(img_height_px / dpi * 914400)

        for shape in slide.shapes:
            if shape.name == "Picture 2":
                box_left, box_top = shape.left, shape.top
                box_width, box_height = shape.width, shape.height

                scale = min(
                    box_width / img_width if img_width > box_width else 1.0,
                    box_height / img_height if img_height > box_height else 1.0
                )
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)

                final_left = box_left + int((box_width - final_width) / 2)
                final_top = box_top + int((box_height - final_height) / 2)

                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(temp_img_path, final_left, final_top, width=final_width, height=final_height)
                break

    # Slide 8 (FOUR images)
    slide = prs.slides[7]
    slide8_img_configs = [
    ("slide_8_first",  "Picture 12", "temp_slide_8_first.jpg"),
    ("slide_8_second", "Picture 13", "temp_slide_8_second.jpg"),
    ("slide_8_third",  "Picture 16", "temp_slide_8_third.jpg"),
    ("slide_8_fourth", "Picture 17", "temp_slide_8_fourth.jpg"),
    ]

    for img_key, shape_name, temp_path in slide8_img_configs:
        if images and img_key in images and images[img_key] is not None:
            img_bytes = images[img_key].read()
            with open(temp_path, "wb") as f:
                f.write(img_bytes)

        img = Image.open(temp_path)
        img_width_px, img_height_px = img.size

        dpi = 96
        img_width = int(img_width_px / dpi * 914400)
        img_height = int(img_height_px / dpi * 914400)

        for shape in slide.shapes:
            if shape.name == shape_name:
                box_left, box_top = shape.left, shape.top
                box_width, box_height = shape.width, shape.height

                scale = min(
                    box_width / img_width if img_width > box_width else 1.0,
                    box_height / img_height if img_height > box_height else 1.0
                )
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)

                final_left = box_left + int((box_width - final_width) / 2)
                final_top = box_top + int((box_height - final_height) / 2)

                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(temp_path, final_left, final_top, width=final_width, height=final_height)
                break


    #slide 11 (Four images)
    slide = prs.slides[10]
    slide11_img_configs = [
    ("slide_11_first", "Picture 26", "temp_slide_11_first.jpg"),
    ("slide_11_second", "Picture 15", "temp_slide_11_second.jpg"),
    ("slide_11_third", "Picture 27", "temp_slide_11_third.jpg"),
    ("slide_11_fourth", "Picture 31", "temp_slide_11_fourth.jpg")
    ]

    for img_key, shape_name, temp_path in slide11_img_configs:
        if images and img_key in images and images[img_key] is not None:
            img_bytes = images[img_key].read()
            with open(temp_path, "wb") as f:
                f.write(img_bytes)

        img = Image.open(temp_path)
        img_width_px, img_height_px = img.size

        dpi = 96
        img_width = int(img_width_px / dpi * 914400)
        img_height = int(img_height_px / dpi * 914400)

        for shape in slide.shapes:
            if shape.name == shape_name:
                box_left, box_top = shape.left, shape.top
                box_width, box_height = shape.width, shape.height

                scale = min(
                    box_width / img_width if img_width > box_width else 1.0,
                    box_height / img_height if img_height > box_height else 1.0
                )
                final_width = int(img_width * scale)
                final_height = int(img_height * scale)

                final_left = box_left + int((box_width - final_width) / 2)
                final_top = box_top + int((box_height - final_height) / 2)

                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(temp_path, final_left, final_top, width=final_width, height=final_height)
                break
    # Social Posts & Stories
    social_posts_value = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Number of Posts With Stories":
                social_posts_value = row["Unnamed: 11"]
                break



    organic_views_impressions = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Organic (Views)":
                organic_views_impressions = row["Unnamed: 11"]
                break



    organic_reach_impressions = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Organic (Reach)":
                organic_reach_impressions = row["Unnamed: 11"]
                break

    impressions_paid = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Paid":
                impressions_paid = row["Unnamed: 11"]
                break



    # Engagements
    engagements_value = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Engagements":
                engagements_value = row["Unnamed: 11"]
                break

    # Impressions
    impressions_value = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            cell_value = str(row["Organic & Total"]).strip()
            if cell_value in ("Total", "Total Impressions"):
                impressions_value = row["Unnamed: 11"]
                break

    # Engagement Rate
    engagement_rate_value = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Program ER":
                engagement_rate_value = row["Unnamed: 11"]
                engagement_rate_value = float(engagement_rate_value) * 100
                engagement_rate_value = str(engagement_rate_value)
                if engagement_rate_value.startswith("0."):
                    engagement_rate_value = engagement_rate_value[1:]
                dot_idx = engagement_rate_value.find(".")
                if dot_idx != -1:
                    engagement_rate_value = engagement_rate_value[:dot_idx + 3]
                break

    # Engagements & Impressions % INCREASE
    engagements_increase = ""
    impressions_increase = ""
    try:
        engagement_val = excel_df.at[5, "Unnamed: 15"]
        impression_val = excel_df.at[4, "Unnamed: 15"]
        if pd.notna(engagement_val):
            engagements_increase = f"{float(engagement_val) * 100:.1f}%"
        if pd.notna(impression_val):
            impressions_increase = f"{float(impression_val) * 100:.1f}%"
    except Exception as e:
        print("could not find value")

    print("Engagements % increase:", engagements_increase)
    print("Impressions % increase:", impressions_increase)

    # Organic Likes
    organic_likes = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Likes":
                organic_likes = row["Unnamed: 11"]
                break

    # Organic Comments
    organic_comments = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Comments":
                organic_comments = row["Unnamed: 11"]
                break

    # Organic Shares
    organic_shares = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Shares":
                organic_shares = row["Unnamed: 11"]
                break

    # Organic Saves
    organic_saves = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Saves":
                organic_saves = row["Unnamed: 11"]
                break
    
    paid_engagements = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Paid Engagements":
                paid_engagements = row["Unnamed: 11"]
                break
    
    paid_likes = ""
    if "Unnamed: 14" in excel_df.columns and "Dates" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 14"]).strip() == "Reactions":
                paid_likes = row["Dates"]
                break

    paid_comments = ""
    if "Unnamed: 14" in excel_df.columns and "Dates" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 14"]).strip() == "Comments":
                paid_comments = row["Dates"]
                break
    
    paid_shares = ""
    if "Unnamed: 14" in excel_df.columns and "Dates" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 14"]).strip() == "Shares":
                paid_shares = row["Dates"]
                break

    
    influencer_count = ""
    if "Dates" in excel_df.columns and "Unnamed: 14" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Dates"]).strip() == "Influencers":
                influencer_count = row["Unnamed: 14"]
                break
    
    diversity_value = ""
    diversity_col = None

# Normalize column names to string and strip whitespace
    for col in excel_df.columns:
        if str(col).strip().lower() == "diversity":
            diversity_col = col
            break

    if diversity_col is not None:
    # Try to get the first non-empty, non-nan value under the Diversity column
        for val in excel_df[diversity_col]:
            if pd.notna(val) and str(val).strip() != "":
                diversity_value = str(val).strip()
                break
    


    paid_saves = ""
    if "Unnamed: 14" in excel_df.columns and "Dates" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 14"]).strip() == "Saves":
                paid_saves = row["Dates"]
                break 
    
    paid_threesec = ""
    if "Unnamed: 14" in excel_df.columns and "Dates" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 14"]).strip() == "3 sec vid views":
                paid_threesec = row["Dates"]
                break
    
    ##total_post_engagements = (
    ##int(organic_likes) + int(organic_comments) + int(organic_shares) + int(organic_saves)
    ##+ int(paid_likes) + int(paid_comments) + int(paid_shares) + int(paid_saves) + int(paid_threesec)
##)
    
    story_engagements = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Story Engagements":
                story_engagements = row["Unnamed: 11"]
                break

    total_engagements = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "Total Engagements":
                total_engagements = row["Unnamed: 11"]
                break

    cpe = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "CPE":
                cpe = row["Unnamed: 17"]
                break
    
    cpc = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "CPC":
                cpc = row["Unnamed: 17"]
                break

    ctr = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "CTR":
                ctr = row["Unnamed: 17"]
                break

    cpm = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "CPM":
                cpm = row["Unnamed: 17"]
                break
    

    thruplays = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "ThruPlays":
                thruplays = row["Unnamed: 17"]
                break

    p25 = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "0.25":
                p25 = row["Unnamed: 17"]
                break

    p50 = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "0.5":
             p50 = row["Unnamed: 17"]
             break

    p75 = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Unnamed: 18"]).strip() == "0.75":
             p75 = row["Unnamed: 17"]
             break

    p100 = ""
    if "Unnamed: 18" in excel_df.columns and "Unnamed: 17" in excel_df.columns:
        for _, row in excel_df.iterrows():
         if str(row["Unnamed: 18"]).strip() == "1":
            p100 = row["Unnamed: 17"]
            break
    
    c2c_transfer = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "C2C Transfers":
                c2c_transfer = row["Unnamed: 11"]
                break
    
    c2c_value = ""
    if "Organic & Total" in excel_df.columns and "Unnamed: 11" in excel_df.columns:
        for _, row in excel_df.iterrows():
            if str(row["Organic & Total"]).strip() == "C2C Value": 
                c2c_value = row["Unnamed: 11"]
         
    
    # Fill TextBox 2 (Proposed Metrics)
    slide = prs.slides[3]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 2":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if "Proposed Influencers" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(metrics.get("Influencers", "")))
                elif "Proposed Engagements" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(metrics.get("Engagements", "")))
                elif "Proposed Impressions" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(metrics.get("Impressions", "")))
                

            

    # Fill TextBox 15 (Program Overview)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 15":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if "Influencers" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(influencer_count))
                elif "Diversity Rate" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(diversity_value))
                # Social Posts & Stories
                elif "Social Posts & Stories" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(social_posts_value))
                # Engagement Rate
                elif "Engagement Rate" in text:
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(engagement_rate_value))
                # Engagements (main)
                elif "Engagements" in text and "#% increase" in text:
                    main_done = False
                    percent_done = False
                    for run in para.runs:
                        if "#" in run.text and not main_done and "% increase" not in run.text:
                            run.text = run.text.replace("#", str(engagements_value), 1)
                            main_done = True
                        if "#% increase" in run.text and not percent_done:
                            run.text = run.text.replace("#", str(engagements_increase), 1)
                            percent_done = True
                # Impressions (main)
                elif "Impressions" in text and "#% increase" in text:
                    main_done = False
                    percent_done = False
                    for run in para.runs:
                        if "#" in run.text and not main_done and "% increase" not in run.text:
                            run.text = run.text.replace("#", str(impressions_value), 1)
                            main_done = True
                        if "#% increase" in run.text and not percent_done:
                            run.text = run.text.replace("#", str(impressions_increase), 1)
                            percent_done = True

    # Fill TextBox 19 (Slide 9 vertical fields)
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 19":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "10" in run.text and "K" in run.text:
                        run.text = run.text.replace("10", str(organic_likes))
                        run.text = run.text.replace("K", "")
                    if "20" in run.text and "K" in run.text:
                        run.text = run.text.replace("20", str(organic_comments))
                        run.text = run.text.replace("K", "")
                    if "30" in run.text and "K" in run.text:
                        run.text = run.text.replace("30", str(organic_shares))
                        run.text = run.text.replace("K", "")
                    if "40" in run.text and "K" in run.text:
                        run.text = run.text.replace("40", str(organic_saves))
                        run.text = run.text.replace("K", "")

# Paid – TextBox 11
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 11":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "10" in run.text and "K" in run.text:
                        run.text = run.text.replace("10", str(paid_likes))
                        run.text = run.text.replace("K", "")
                    if "20" in run.text and "K" in run.text:
                        run.text = run.text.replace("20", str(paid_comments))
                        run.text = run.text.replace("K", "")
                    if "30" in run.text and "K" in run.text:
                        run.text = run.text.replace("30", str(paid_shares))
                        run.text = run.text.replace("K", "")
                    if "40" in run.text and "K" in run.text:
                        run.text = run.text.replace("40", str(paid_saves))
                        run.text = run.text.replace("K", "")
                    if "##" in run.text and "K" in run.text:
                        run.text = run.text.replace("##", str(paid_threesec))
                        run.text = run.text.replace("K", "")
    
    #slide 9 last boxes
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 34":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "100" in run.text and "K" in run.text:
                        ##run.text = run.text.replace("100", str(total_post_engagements))
                        run.text = run.text.replace("K", "")
                    if "200" in run.text and "K" in run.text:
                        run.text = run.text.replace("200", str(story_engagements))
                        run.text = run.text.replace("K", "")
        
        elif shape.has_text_frame and shape.name == "TextBox 18":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "222" in run.text and "K" in run.text:
                        run.text = run.text.replace("222", str(total_engagements))
                        run.text = run.text.replace("K", "")


    #slide 10
    slide = prs.slides[9]
    for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "TextBox 18":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(organic_reach_impressions))

            elif shape.has_text_frame and shape.name == "TextBox 19":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(impressions_paid))
            
            elif shape.has_text_frame and shape.name == "TextBox 21":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(organic_views_impressions))
           
            elif shape.has_text_frame and shape.name == "TextBox 29":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(impressions_value))
    

    #slide 11
    slide = prs.slides[10]
    for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "TextBox 18":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(organic_reach_impressions))

            elif shape.has_text_frame and shape.name == "TextBox 19":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(impressions_paid))
            
            elif shape.has_text_frame and shape.name == "TextBox 21":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(organic_views_impressions))
           
            elif shape.has_text_frame and shape.name == "TextBox 29":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(impressions_value))
    
    
    #slide 6 text
    slide = prs.slides[5]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 9":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "influencerhandle" in run.text:
                        run.text = run.text.replace("influencerhandle", handle_slide_6)

    #slide 7 text
    slide = prs.slides[6]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 6":
            if "Organic" in shape.text:
                hashtag_values = [slide_7_likes, slide_7_comments, slide_7_views, slide_7_reach]
                value_index = 0
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                # Replace influencer handle
                        if "influencerhandle" in run.text:
                            run.text = run.text.replace("influencerhandle", handle_slide_7_left)
                # Replace hashtags one by one in order
                        if "#" in run.text and value_index < len(hashtag_values):
                            run.text = run.text.replace("#", str(hashtag_values[value_index]))
                            value_index += 1
            elif "Paid" in shape.text:
                hashtag_values_paid = [slide_7_engagements, slide_7_impressions]
                value_index = 0
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "influencerhandle" in run.text:
                            run.text = run.text.replace("influencerhandle", handle_slide_7_right)
                        if "#" in run.text and value_index < len(hashtag_values_paid):
                            run.text = run.text.replace("#", str(hashtag_values_paid[value_index]))
                            value_index += 1

    
    #slide 12


    slide = prs.slides[11]
    label_to_value = {
    "CPE": str(cpe),
    "CPC": str(cpc),
    "CTR": str(ctr),
    "CPM": str(cpm),
    }

    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 6":
        # Find the right TextBox 6 by its content
            if any(label in para.text for label in label_to_value.keys() for para in shape.text_frame.paragraphs):
                for para in shape.text_frame.paragraphs:
                    for label, value in label_to_value.items():
                        if label in para.text:
                        # Prepend value and a space to the first run
                            if para.runs:
                                para.runs[0].text = f"{value} " + para.runs[0].text
                            break  # Only update once per paragraph
                break


# For ThruPlays
    slide = prs.slides[11]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 6":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "# ThruPlays":
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(thruplays))
# ...existing code...



    slide = prs.slides[11]
    for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "TextBox 13":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(p25))

            elif shape.has_text_frame and shape.name == "TextBox 11":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(p50))
            
            elif shape.has_text_frame and shape.name == "TextBox 3":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(p75))
           
            elif shape.has_text_frame and shape.name == "TextBox 26":
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    for run in para.runs:
                        if "#" in run.text:
                            run.text = run.text.replace("#", str(p100))



#slide 4 program goals
    slide = prs.slides[3]  # Or whatever index slide 4 actually is (Python is 0-based)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 10":
            for para in shape.text_frame.paragraphs:
            # Bullet 1
                if para.text.strip() == "Create excitement and promote (brand) products available at (retailer).":
                    if para.runs:
                        para.runs[0].text = str(slide_4_bullet1)
                        for run in para.runs[1:]:
                            run.text = ""
            # Bullet 2
                if para.text.strip() == "Encourage shoppers to purchase the (brand and products)…":
                    if para.runs:
                        para.runs[0].text = str(slide_4_bullet2)
                        for run in para.runs[1:]:
                            run.text = ""


# Slide 9 text (replace entire line if it matches the placeholder)
    slide = prs.slides[8]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 2":
            for para in shape.text_frame.paragraphs:
            # Check if the entire paragraph matches the placeholder
                if para.text.strip() == "Total engagements outperformed proposed estimated engagements (#) by #%.":
                # Replace all runs in this paragraph with the new text, preserving formatting of the first run
                    if para.runs:
                        para.runs[0].text = str(text_slide_9)
                    # Clear out any extra runs (if present)
                        for run in para.runs[1:]:
                            run.text = ""
    
    #slide 1 input
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 5":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "January 1, 2025 – February 1, 2025":
                    if para.runs:
                        para.runs[0].text = str(date_slide_1)
                        for run in para.runs[1:]:
                            run.text = ""
        if shape.has_text_frame and shape.name == "TextBox 6":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "#CampaignHashtag":
                    if para.runs:
                        para.runs[0].text = str(hashtag_slide_1)
                        for run in para.runs[1:]:
                            run.text = ""

#slide 2 input
    slide = prs.slides[1]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 7":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "January 1, 2025 – February 1, 2025":
                    if para.runs:
                        para.runs[0].text = str(date_slide_2)
                        for run in para.runs[1:]:
                            run.text = ""
        if shape.has_text_frame and shape.name == "TextBox 8":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "#CampaignHashtag":
                    if para.runs:
                        para.runs[0].text = str(hashtag_slide_2)
                        for run in para.runs[1:]:
                            run.text = ""

#slide 3 input
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 7":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "January 1, 2025 – February 1, 2025":
                    if para.runs:
                        para.runs[0].text = str(date_slide_3)
                        for run in para.runs[1:]:
                            run.text = ""
        if shape.has_text_frame and shape.name == "TextBox 8":
            for para in shape.text_frame.paragraphs:
                if para.text.strip() == "#CampaignHashtag":
                    if para.runs:
                        para.runs[0].text = str(hashtag_slide_3)
                        for run in para.runs[1:]:
                            run.text = ""






    influencer_boxes = text_inputs.get("influencer_boxes", {})
    influencer_boxestwo = text_inputs.get("influencer_boxestwo", [])

#slide 5 text inputs
    slide = prs.slides[4]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name in influencer_boxes:
            replacements = influencer_boxes[shape.name]
        # Combine city/state if needed
            city_state = f"{replacements.get('City','')}, {replacements.get('State','')}".strip(", ")
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                # Replace handle
                    if "influencerhandle" in run.text:
                        run.text = run.text.replace("influencerhandle", replacements.get("influencerhandle", ""))
                # Replace reach
                    if "##" in run.text:
                        run.text = run.text.replace("##", replacements.get("##", ""))
                # Replace city, state
                    if "City, State" in run.text:
                        run.text = run.text.replace("City, State", city_state)
                # Replace verbatim (with or without quotes)
                    if "Verbatim" in run.text:
                        run.text = run.text.replace("Verbatim", replacements.get("Verbatim", ''))
                    elif "Verbatim" in run.text:
                        run.text = run.text.replace("Verbatim", replacements.get("Verbatim", ""))

#slide 8 text inputs
    slide = prs.slides[7]  # Use your slide index
    box_index = 0
    metric_keys = ["# Likes", "# Comments", "# Views", "# Social Reach"]

    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 6":
            if box_index < len(influencer_boxestwo):
                replacements = influencer_boxestwo[box_index]
                paras = shape.text_frame.paragraphs

            # First paragraph: influencerhandle
                if len(paras) > 0:
                    for run in paras[0].runs:
                        if "influencerhandle" in run.text:
                            run.text = run.text.replace("influencerhandle", replacements.get("influencerhandle", ""))

            # Next paragraphs: metrics
                for i, key in enumerate(metric_keys):
                    para_idx = i + 1  # starts from second paragraph
                    if para_idx < len(paras):
                        for run in paras[para_idx].runs:
                            if "#" in run.text:
                                run.text = run.text.replace("#", replacements.get(key, ""))

            box_index += 1
            

#slide 13 data
    slide = prs.slides[12]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 5":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "10.6K" in run.text:
                        run.text = run.text.replace("10.6K", str(c2c_transfer))
        if shape.has_text_frame and shape.name == "TextBox 21":
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                for run in para.runs:
                    if "27K" in run.text:
                        run.text = run.text.replace("27K", str(c2c_value))


#slide 13 text
    slide = prs.slides[12]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 3":
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "00/00/00 – 00/00/00" in run.text:
                        run.text = run.text.replace("00/00/00 – 00/00/00", str(text_slide_13))


#slide 15 text
    slide = prs.slides[14]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 5":
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "On a scale from 1 to 10, 10 being the most likely, how likely would you be able to recommend Ticket to Ride/Ticket to Ride: San Francisco to family and friends?" in run.text:
                        run.text = run.text.replace("On a scale from 1 to 10, 10 being the most likely, how likely would you be able to recommend Ticket to Ride/Ticket to Ride: San Francisco to family and friends?", str(text_slide_15))


#slide 16 text
    slide = prs.slides[15]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "TextBox 5":
            found = False  # Track if we've already inserted the user text
            paras_to_remove = []
            for i, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if (
                text == "What were your favorite parts of the game night"
                or text == "Playing Ticket to Ride: San Francisco?"
                ):
                    if not found:
                        if para.runs:
                            para.runs[0].text = str(text_slide_16)
                            for run in para.runs[1:]:
                                run.text = ""
                        found = True  # Only insert user entry once
                    else:
                    # Mark this extra placeholder paragraph for removal
                        paras_to_remove.append(i)
        # Remove the second matched paragraph (after loop, so we don't break indexing)
        # Remove from the end so earlier indices remain correct
            for idx in reversed(paras_to_remove):
                shape.text_frame._element.remove(shape.text_frame.paragraphs[idx]._element)





         
    prs.save(output_path)
                        
      

# ─────────────────────────────────────────────────────────────────────────────
# CLI Entrypoint (optional, for testing)
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Populate PowerPoint deck from Excel data")
    parser.add_argument("input_file", help="CSV or Excel input")
    parser.add_argument("pptx_template", help="PowerPoint template file")
    parser.add_argument("--output", default="recap_deck.pptx", help="Output PPTX file")
    args = parser.parse_args()

    df = load_dataframe(args.input_file)
    populate_pptx_from_excel(df, args.pptx_template, args.output)
    print(f"Wrote {args.output}")
